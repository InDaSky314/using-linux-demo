"""Generate the Using Linux — Manually and Automated demo deck.

Audience-facing slides: opens with Linux foundations (kernel, multi-user, FHS, permission
bits), walks the by-hand phase on canvas, then the ansible automation phase, and closes
with where ansible sits in the wider IaC stack.

Run from anywhere — the .pptx is written next to this script:
    python3 build_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ── Theme palette (sampled from Linux_From_The_Ground_Up.pptx) ─────────────
BG          = RGBColor(0x0D, 0x11, 0x17)
BAND        = RGBColor(0x13, 0x1F, 0x2E)
CODE_BG     = RGBColor(0x0A, 0x0E, 0x14)
CARD        = RGBColor(0x1A, 0x2B, 0x3C)
ALT_ROW     = RGBColor(0x12, 0x1E, 0x2C)
GREEN       = RGBColor(0x39, 0xD3, 0x53)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
RED         = RGBColor(0xE5, 0x57, 0x5A)
TEXT        = RGBColor(0xC8, 0xD6, 0xE5)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
MUTED       = RGBColor(0x6B, 0x7E, 0x8C)

CALIBRI = "Calibri"
CONSOLAS = "Consolas"
IMPACT = "Impact"

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Using_Linux_Manually_and_Automated.pptx")

# ── Setup ──────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if not line:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, font=CALIBRI, size=12, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_multiline(slide, x, y, w, h, lines, font=CALIBRI, size=11, color=TEXT,
                  bold=False, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = text
        r.font.name = opts.get("font", font)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.color.rgb = opts.get("color", color)
    return tb


def slide_bg(slide):
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, BG)


def header_band(slide, title, source=None):
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.85), BAND)
    add_text(slide, Inches(0.20), Inches(0.10), Inches(12.0), Inches(0.65),
             title, font=CALIBRI, size=24, bold=True, color=WHITE)
    add_rect(slide, 0, Inches(7.15), prs.slide_width, Inches(0.35), BAND)
    if source:
        add_text(slide, Inches(0.30), Inches(7.18), Inches(12.5), Inches(0.28),
                 source, font=CONSOLAS, size=8, color=MUTED)


def card(slide, x, y, w, h, title, accent=GREEN, items=None, body=None,
         body_size=11, item_size=11):
    add_rect(slide, x, y, w, h, CARD)
    add_rect(slide, x, y, Emu(54864), h, accent)
    add_text(slide, x + Inches(0.15), y + Inches(0.05),
             w - Inches(0.30), Inches(0.35),
             title, font=CALIBRI, size=13, bold=True, color=WHITE)
    if items:
        cur = y + Inches(0.42)
        line_h = Inches(0.36)
        for it in items:
            add_text(slide, x + Inches(0.15), cur, w - Inches(0.30), line_h,
                     it, font=CALIBRI, size=item_size, color=TEXT)
            cur += line_h
    if body:
        add_text(slide, x + Inches(0.15), y + Inches(0.42),
                 w - Inches(0.30), h - Inches(0.5),
                 body, font=CALIBRI, size=body_size, color=TEXT)


def code_block(slide, x, y, w, h, header, lines, accent=GREEN):
    add_rect(slide, x, y, w, h, CODE_BG)
    add_rect(slide, x, y, Emu(45720), h, accent)
    add_rect(slide, x, y, w, Inches(0.30), BAND)
    add_text(slide, x + Inches(0.10), y + Inches(0.02), w - Inches(0.20), Inches(0.26),
             header, font=CONSOLAS, size=9, color=MUTED)
    cur = y + Inches(0.34)
    line_h = Inches(0.24)
    for ln in lines:
        if isinstance(ln, tuple):
            text, color = ln
        else:
            text, color = ln, TEXT
        add_text(slide, x + Inches(0.12), cur, w - Inches(0.24), line_h,
                 text, font=CONSOLAS, size=10, color=color)
        cur += line_h


def section_divider(prs, section_label, big_title, subtitle):
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    add_rect(s, 0, 0, Inches(0.30), prs.slide_height, GREEN)
    add_text(s, Inches(0.70), Inches(1.50), Inches(10), Inches(0.60),
             section_label, font=CONSOLAS, size=14, bold=True, color=GREEN)
    add_text(s, Inches(0.70), Inches(2.10), Inches(11), Inches(1.50),
             big_title, font=CALIBRI, size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.70), Inches(3.80), Inches(11), Inches(0.80),
             subtitle, font=CALIBRI, size=20, color=TEXT)
    add_rect(s, 0, Inches(7.10), prs.slide_width, Inches(0.40), BAND)
    add_text(s, Inches(0.30), Inches(7.12), Inches(12), Inches(0.30),
             "USING LINUX  ●  MANUALLY AND AUTOMATED  ●  CAREER PRACTICUM",
             font=CONSOLAS, size=9, color=MUTED)
    return s


def content_slide(title, source):
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    header_band(s, title, source)
    return s


# ─── SLIDE 1: TITLE ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
for x in range(0, 14):
    add_rect(s, Inches(x), 0, Emu(9144), prs.slide_height, BAND)
for y in range(0, 8):
    add_rect(s, 0, Inches(y), prs.slide_width, Emu(9144), BAND)
add_rect(s, 0, 0, Inches(0.50), prs.slide_height, GREEN)
add_text(s, Inches(0.70), Inches(0.50), Inches(12), Inches(2.50),
         "USING LINUX", font=IMPACT, size=110, bold=True, color=WHITE)
add_rect(s, Inches(0.70), Inches(2.85), Inches(9), Emu(54864), GREEN)
add_text(s, Inches(0.70), Inches(2.95), Inches(12), Inches(0.70),
         "MANUALLY AND AUTOMATED",
         font=CONSOLAS, size=28, color=TEXT)
add_text(s, Inches(0.70), Inches(3.65), Inches(12), Inches(0.45),
         "Kernel  ●  Users & permissions  ●  Filesystem  ●  Ansible primitives  ●  Idempotency",
         font=CALIBRI, size=14, color=MUTED)
add_text(s, Inches(0.70), Inches(4.20), Inches(10), Inches(0.40),
         "Career Practicum — Live Demonstration",
         font=CALIBRI, size=14, color=TEXT)
for i, label in enumerate(["bash", "useradd", "chmod", "/srv", "playbook", "idempotent", "Piwigo"]):
    add_rect(s, Inches(0.70 + i*1.74), Inches(4.75), Inches(1.60), Inches(0.32), BAND)
    add_text(s, Inches(0.75 + i*1.74), Inches(4.77), Inches(1.50), Inches(0.28),
             label, font=CONSOLAS, size=10, color=GREEN)
add_rect(s, 0, Inches(7.10), prs.slide_width, Inches(0.40), BAND)
add_text(s, Inches(0.30), Inches(7.13), Inches(12), Inches(0.28),
         "Hosts: brush 192.168.50.5 (Ubuntu 24.04)  ●  canvas 192.168.50.26 (Ubuntu 24.04)  ●  Kernel 6.8",
         font=CONSOLAS, size=9, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SECTION 01 — LINUX FOUNDATIONS
# ═══════════════════════════════════════════════════════════════════════════
section_divider(prs, "SECTION 01",
                "LINUX FOUNDATIONS",
                "The kernel, the multi-user model, the filesystem, and the permission bits everything else builds on")


# ─── SLIDE 3: WHAT IS LINUX? ───────────────────────────────────────────────
s = content_slide("What Linux Is — Kernel, Userland, Distribution",
                  "Source: kernel.org  |  Linus Torvalds, comp.os.minix Aug 1991  |  gnu.org")

add_text(s, Inches(0.30), Inches(0.95), Inches(12.7), Inches(0.95),
         "Linux is the kernel — the program that talks to hardware and mediates everything else. "
         "The operating system in everyday use is the kernel plus the GNU userland (bash, coreutils, "
         "systemd) bundled by a distribution. Today's demo runs on Ubuntu 24.04 — a Debian-derived "
         "distribution with the Linux 6.8 kernel.",
         font=CALIBRI, size=12, color=TEXT)

card(s, Inches(0.30), Inches(2.05), Inches(4.20), Inches(3.10),
     "Kernel (Linux proper)", accent=GREEN,
     items=["▸  CPU scheduling, memory mgmt",
            "▸  Device drivers, network stack",
            "▸  Provides system calls",
            "▸  Enforces every permission check",
            "▸  Monolithic, ~30M lines of C",
            "▸  GPLv2 — free + open source"], item_size=12)

card(s, Inches(4.65), Inches(2.05), Inches(4.20), Inches(3.10),
     "GNU + Userland", accent=CYAN,
     items=["▸  Shell — bash, zsh",
            "▸  Coreutils — ls, cat, chmod",
            "▸  useradd, groupadd, install",
            "▸  Init / service mgr — systemd",
            "▸  Networking — iproute2, ssh",
            "▸  Package mgr — apt, dnf, pacman"], item_size=12)

card(s, Inches(9.00), Inches(2.05), Inches(4.00), Inches(3.10),
     "Distribution", accent=ORANGE,
     items=["▸  Ubuntu (our demo OS)",
            "▸  Debian, Fedora, RHEL, Arch",
            "▸  Bundles kernel + userland",
            "▸  Adds package manager",
            "▸  Release eng + security updates",
            "▸  ~600 actively-maintained distros"], item_size=12)

code_block(s, Inches(0.30), Inches(5.30), Inches(12.70), Inches(1.75),
           "$ verifying the kernel and distro on canvas",
           ["$ uname -r",
            "6.8.0-60-generic",
            "$ uname -a",
            "Linux canvas 6.8.0-60-generic #63-Ubuntu SMP PREEMPT_DYNAMIC x86_64 GNU/Linux",
            "$ lsb_release -d",
            "Description:  Ubuntu 24.04.1 LTS"])


# ─── SLIDE 4: MULTI-USER OS ───────────────────────────────────────────────
s = content_slide("Multi-User from Day One — Identity Is the Permission System",
                  "Source: passwd(5)  |  group(5)  |  id(1)  |  /etc/passwd, /etc/shadow, /etc/group")

add_text(s, Inches(0.30), Inches(0.95), Inches(12.7), Inches(0.95),
         "Linux has been multi-user since the first release. Every running process belongs to a user "
         "and one or more groups; every file on disk has an owning user and group. The kernel checks "
         "every read, write, and execute against those identities — there is no untracked access.",
         font=CALIBRI, size=12, color=TEXT)

card(s, Inches(0.30), Inches(2.05), Inches(6.30), Inches(2.10),
     "/etc/passwd  —  one line per user", accent=GREEN,
     items=["username : x : UID : GID : GECOS : home : shell",
            "",
            "rembrandt:x:1001:1001::/home/rembrandt:/bin/bash",
            "(x = password lives in /etc/shadow, root-only)"], item_size=12)

card(s, Inches(6.70), Inches(2.05), Inches(6.30), Inches(2.10),
     "/etc/group  —  one line per group", accent=CYAN,
     items=["groupname : x : GID : member,member,member",
            "",
            "painters:x:1001:rembrandt,klimt,dali,kahlo",
            "(comma-separated supplementary group members)"], item_size=12)

code_block(s, Inches(0.30), Inches(4.30), Inches(12.70), Inches(2.75),
           "$ identity in action",
           ["$ id rembrandt",
            "uid=1001(rembrandt) gid=1001(painters) groups=1001(painters)",
            "",
            "$ ls -l /srv/atelier/rembrandt/sketches.txt",
            "-rw-------  1 rembrandt painters  35 May  6 09:00 sketches.txt",
            "                ↑ owner    ↑ group   ↑ kernel checks identity on every open(2)",
            "",
            "$ ps -eo uid,user,comm | head -3",
            "    0 root             systemd",
            "  101 systemd-network  systemd-network"])


# ─── SLIDE 5: FILESYSTEM HIERARCHY ─────────────────────────────────────────
s = content_slide("The Filesystem — Where Things Belong (FHS)",
                  "Source: pathname.com/fhs  |  hier(7) man page  |  Filesystem Hierarchy Standard 3.0")

add_text(s, Inches(0.30), Inches(0.95), Inches(12.7), Inches(0.55),
         "Linux uses a single rooted tree — everything hangs off /. The Filesystem Hierarchy Standard "
         "gives every directory a defined purpose, so admins on any distro know where to look.",
         font=CALIBRI, size=12, color=TEXT)

card(s, Inches(0.30), Inches(1.70), Inches(6.30), Inches(3.40),
     "Directories you'll see all the time", accent=GREEN,
     items=["▸  /etc      — system configuration files",
            "▸  /home     — regular users' home dirs",
            "▸  /var      — variable data: logs, mail, spool",
            "▸  /tmp      — scratch, world-writable, sticky",
            "▸  /usr      — read-only userland (binaries, libs)",
            "▸  /opt      — third-party software"], item_size=12)

card(s, Inches(6.70), Inches(1.70), Inches(6.30), Inches(3.40),
     "Where the demo creates things", accent=ORANGE,
     items=["▸  /srv      — service-data root  (FHS reserved)",
            "▸  /srv/atelier  — Phase 1 (manual) studio",
            "▸  /srv/studio   — Phase 2 (ansible) studio",
            "▸  /var/www/piwigo  — Apache document root",
            "▸  /etc/apache2/sites-available/piwigo.conf",
            "▸  /home/<artist>  — created by useradd -m"], item_size=12)

code_block(s, Inches(0.30), Inches(5.30), Inches(12.70), Inches(1.75),
           "$ poking around the tree",
           ["$ ls /",
            "bin   boot  dev  etc  home  lib  media  mnt  opt  proc  root  run  sbin  srv  sys  tmp  usr  var",
            "",
            "$ stat -c '%n  %U:%G  %a' /etc /home /srv /tmp",
            "/etc  root:root  755     /home  root:root  755",
            "/srv  root:root  755     /tmp   root:root  1777"])


# ─── SLIDE 6: PERMISSION BITS ──────────────────────────────────────────────
s = content_slide("Permission Bits — The Model the Kernel Enforces",
                  "Source: chmod(1)  |  inode(7)  |  ls(1) -l output format")

# big visual at top
add_rect(s, Inches(0.30), Inches(0.95), Inches(12.70), Inches(0.85), CODE_BG)
add_text(s, Inches(0.40), Inches(0.97), Inches(12.5), Inches(0.50),
         "-     rwx        r-x        r--",
         font=CONSOLAS, size=28, bold=True, color=GREEN)
add_text(s, Inches(0.40), Inches(1.45), Inches(12.5), Inches(0.30),
         "type   owner     group      others",
         font=CONSOLAS, size=12, color=MUTED)

# permission table
def table_row(y, cells, fill, bold=False):
    widths = [1.50, 1.10, 1.10, 1.10, 8.00]
    x = 0.30
    for w, text in zip(widths, cells):
        add_rect(s, Inches(x), Inches(y), Inches(w - 0.02), Inches(0.42), fill)
        add_text(s, Inches(x + 0.10), Inches(y + 0.06), Inches(w - 0.20), Inches(0.32),
                 text, font=CONSOLAS, size=11, bold=bold, color=WHITE if bold else TEXT)
        x += w

table_row(2.00, ["Permission", "Symbol", "Octal", "Binary", "Meaning"], BAND, bold=True)
rows = [
    ("Read",    "r",   "4",     "100",  "View file contents / list directory"),
    ("Write",   "w",   "2",     "010",  "Modify file / create + delete in directory"),
    ("Execute", "x",   "1",     "001",  "Run as program / cd into directory"),
    ("None",    "-",   "0",     "000",  "Denied"),
    ("SUID",    "s(u)", "4000", "—",    "Run as file owner (e.g. /usr/bin/sudo)"),
    ("SGID",    "s(g)", "2000", "—",    "Run as file group / dirs: inherit group ← used in demo"),
    ("Sticky",  "t",   "1000",  "—",    "Only owner can delete (e.g. /tmp)"),
]
y = 2.42
for i, row in enumerate(rows):
    table_row(y, list(row), ALT_ROW if i % 2 == 0 else CARD)
    y += 0.40

code_block(s, Inches(0.30), Inches(5.55), Inches(12.70), Inches(1.50),
           "$ chmod by example",
           ["chmod 644 portfolio.txt        # rw-r--r--   owner rw, others r",
            "chmod 600 ~/.ssh/id_ed25519    # rw-------   ssh REQUIRES this",
            "chmod 755 script.sh            # rwxr-xr-x   executable to all",
            "chmod 2775 /srv/atelier        # rwxrwsr-x   setgid: new files inherit group ← demo uses"])


# ═══════════════════════════════════════════════════════════════════════════
# SECTION 02 — PHASE 1: BY HAND
# ═══════════════════════════════════════════════════════════════════════════
section_divider(prs, "SECTION 02",
                "PHASE 1 — BY HAND",
                "Six users, one shared studio, and a kernel that enforces every bit we just learned about")


# ─── SLIDE 8: THE SETUP ────────────────────────────────────────────────────
s = content_slide("The Setup — Two Hosts, One Job",
                  "Source: /root/ansible/inventory.ini on brush  |  ansible-playbook(1)")

add_text(s, Inches(0.30), Inches(0.95), Inches(12.7), Inches(0.55),
         "Two Ubuntu VMs on the same subnet. Brush holds the playbooks. Canvas receives every change. "
         "During Phase 2, no one ever logs into canvas directly — that is the entire point of the model.",
         font=CALIBRI, size=13, color=TEXT)

add_rect(s, Inches(0.30), Inches(1.70), Inches(6.20), Inches(2.30), CARD)
add_rect(s, Inches(0.30), Inches(1.70), Emu(54864), Inches(2.30), GREEN)
add_text(s, Inches(0.45), Inches(1.78), Inches(6), Inches(0.40),
         "brush — Ansible Control Node", font=CALIBRI, size=15, bold=True, color=WHITE)
add_text(s, Inches(0.45), Inches(2.20), Inches(6), Inches(0.30),
         "192.168.50.5", font=CONSOLAS, size=14, color=GREEN)
add_text(s, Inches(0.45), Inches(2.55), Inches(6), Inches(1.40),
         "Holds the playbooks under /root/ansible/.\n"
         "Runs ansible-playbook against canvas over SSH.\n"
         "Has artists.yml, piwigo.yml, reset.yml, and a one-host inventory.ini.",
         font=CALIBRI, size=12, color=TEXT)

add_rect(s, Inches(6.80), Inches(1.70), Inches(6.20), Inches(2.30), CARD)
add_rect(s, Inches(6.80), Inches(1.70), Emu(54864), Inches(2.30), CYAN)
add_text(s, Inches(6.95), Inches(1.78), Inches(6), Inches(0.40),
         "canvas — Managed Host", font=CALIBRI, size=15, bold=True, color=WHITE)
add_text(s, Inches(6.95), Inches(2.20), Inches(6), Inches(0.30),
         "192.168.50.26", font=CONSOLAS, size=14, color=CYAN)
add_text(s, Inches(6.95), Inches(2.55), Inches(6), Inches(1.40),
         "Receives every change — users, groups, files, packages.\n"
         "No ansible installed here. Just sshd + python3.\n"
         "Holds two parallel namespaces: painters (manual) and artists (ansible).",
         font=CALIBRI, size=12, color=TEXT)

add_rect(s, Inches(0.30), Inches(4.20), Inches(12.70), Inches(0.50), BAND)
add_text(s, Inches(0.45), Inches(4.27), Inches(12.40), Inches(0.40),
         "ssh root@brush  →  ansible-playbook artists.yml  →  ssh + python3 on canvas  →  /srv/studio appears",
         font=CONSOLAS, size=12, color=GREEN, align=PP_ALIGN.CENTER)

code_block(s, Inches(0.30), Inches(4.95), Inches(12.70), Inches(2.10),
           "$ /root/ansible/inventory.ini",
           ["[canvas]",
            "canvas ansible_host=192.168.50.26 ansible_user=root",
            "",
            "[canvas:vars]",
            "ansible_python_interpreter=/usr/bin/python3"])


# ─── SLIDE 9: PHASE 1.1 — GROUPS & USERS ───────────────────────────────────
s = content_slide("Phase 1.1 — Creating the Group and Users on canvas",
                  "Source: groupadd(8)  |  useradd(8)  |  /etc/passwd, /etc/group format")

add_text(s, Inches(0.30), Inches(0.95), Inches(12.7), Inches(0.55),
         "The painters namespace is created entirely by hand. Six users, one group, all created with bash. "
         "This is the ground-truth version of what Phase 2 will collapse into a single playbook task.",
         font=CALIBRI, size=12, color=TEXT)

code_block(s, Inches(0.30), Inches(1.70), Inches(8.10), Inches(3.40),
           "$ root@canvas  —  manual user/group creation",
           ["# Show the group does not exist yet",
            ("getent group painters         # → no output", MUTED),
            "",
            "# Create the group",
            "groupadd painters",
            ("getent group painters         # → painters:x:1001:", MUTED),
            "",
            "# Six users in one for-loop, all in painters",
            "for u in rembrandt klimt dali kahlo banksy okeeffe; do",
            "    useradd -m -s /bin/bash -g painters \"$u\"",
            "done",
            "",
            ("id rembrandt    # uid=1001 gid=1001(painters) groups=1001(painters)", MUTED)])

card(s, Inches(8.60), Inches(1.70), Inches(4.40), Inches(3.40),
     "Scale matters", accent=GREEN,
     items=["▸  One group, six users, eight",
            "    bash commands.",
            "",
            "▸  At one host, the for-loop is",
            "    fine.",
            "",
            "▸  At 200 hosts, repeating it",
            "    by SSH is not."])

code_block(s, Inches(0.30), Inches(5.30), Inches(12.70), Inches(1.75),
           "$ what the same work looks like in artists.yml",
           ["- name: Create artists group",
            "  ansible.builtin.group: { name: artists, state: present }",
            "",
            "- name: Create artist users",
            "  ansible.builtin.user: { name: \"{{ item }}\", group: artists, create_home: true, shell: /bin/bash }",
            "  loop: [vangogh, picasso, davinci, monet, frida, warhol]"],
           accent=CYAN)


# ─── SLIDE 10: PHASE 1.2/1.3 — PERMISSIONS & LIVE PROOF ────────────────────
s = content_slide("Phase 1.2 / 1.3 — Permissions, Setgid, and a Live Permission-Denied",
                  "Source: chmod(1)  |  install(1)  |  setgid behaviour on directories")

add_rect(s, Inches(0.30), Inches(0.95), Inches(12.70), Inches(0.55), CODE_BG)
add_text(s, Inches(0.40), Inches(0.97), Inches(12.5), Inches(0.50),
         "drwxr-sr-x  root  painters  4096  /srv/atelier",
         font=CONSOLAS, size=22, bold=True, color=GREEN)

card(s, Inches(0.30), Inches(1.60), Inches(6.20), Inches(3.10),
     "What the chmod numbers mean (Phase 1.2)", accent=GREEN,
     items=["▸  2775 on /srv/atelier — setgid + rwxrwxr-x",
            "    new files inherit the painters group",
            "▸  0750 on per-user dirs — owner+group only",
            "▸  0600 on sketches.txt — owner read/write only",
            "▸  0644 on portfolio.txt — world-readable",
            "▸  0664 on shared-masterpiece.txt — group writable"])

card(s, Inches(6.80), Inches(1.60), Inches(6.20), Inches(3.10),
     "The kernel enforces it in real time (Phase 1.3)", accent=ORANGE,
     items=["▸  sudo -u rembrandt cat sketches.txt",
            "    → \"Private sketches by rembrandt.\"",
            "▸  sudo -u klimt cat sketches.txt",
            "    → cat: ...: Permission denied",
            "▸  sudo -u klimt tee -a shared-masterpiece.txt",
            "    → succeeds (mode 664 = group writable)"])

code_block(s, Inches(0.30), Inches(4.85), Inches(12.70), Inches(2.20),
           "$ representative chmod calls from /srv/atelier setup",
           ["chmod 2775 /srv/atelier              # 2 = setgid bit",
            "install -d -o rembrandt -g painters -m 0750 /srv/atelier/rembrandt",
            "chmod 600 /srv/atelier/rembrandt/sketches.txt    # private",
            "chmod 644 /srv/atelier/rembrandt/portfolio.txt   # world-readable",
            "chmod 664 /srv/atelier/shared-masterpiece.txt    # group-writable"])


# ═══════════════════════════════════════════════════════════════════════════
# SECTION 03 — PHASE 2: AUTOMATION WITH ANSIBLE
# ═══════════════════════════════════════════════════════════════════════════
section_divider(prs, "SECTION 03",
                "PHASE 2 — WITH ANSIBLE",
                "The same work, declared once in YAML, run anywhere, safe to re-run")


# ─── SLIDE 12: WHAT IS ANSIBLE? ────────────────────────────────────────────
s = content_slide("What Ansible Is — Agentless, Declarative, Idempotent",
                  "Source: docs.ansible.com  |  github.com/ansible/ansible  |  Red Hat AAP")

add_text(s, Inches(0.30), Inches(0.95), Inches(12.7), Inches(0.95),
         "Ansible is a configuration-management tool: you describe what state you want a host to be in, "
         "and ansible makes it so. It runs on a control node (here, brush) and pushes changes over SSH "
         "to managed hosts (here, canvas). The managed host needs nothing installed beyond sshd + python.",
         font=CALIBRI, size=12, color=TEXT)

card(s, Inches(0.30), Inches(2.05), Inches(3.10), Inches(3.10),
     "Agentless", accent=GREEN,
     items=["▸  No daemon on targets",
            "▸  Push over SSH",
            "▸  Only needs sshd + python3",
            "▸  Easy to bolt onto existing fleets"], item_size=11)

card(s, Inches(3.55), Inches(2.05), Inches(3.10), Inches(3.10),
     "Declarative", accent=CYAN,
     items=["▸  YAML describes desired state",
            "▸  Not a sequence of commands",
            "▸  Modules know how to converge",
            "▸  Reads like documentation"], item_size=11)

card(s, Inches(6.80), Inches(2.05), Inches(3.10), Inches(3.10),
     "Idempotent", accent=ORANGE,
     items=["▸  Modules check before acting",
            "▸  ok: when no change needed",
            "▸  changed: when work was done",
            "▸  Safe to run on a schedule"], item_size=11)

card(s, Inches(10.05), Inches(2.05), Inches(2.95), Inches(3.10),
     "Modular", accent=RED,
     items=["▸  ~5,000 built-in + community",
            "▸  group, user, file, copy, apt",
            "▸  service, mysql_db, k8s, aws_ec2",
            "▸  Roles + collections for reuse"], item_size=11)

code_block(s, Inches(0.30), Inches(5.30), Inches(12.70), Inches(1.75),
           "$ smallest possible ansible — does the target answer",
           ["$ ansible canvas -i inventory.ini -m ping",
            "canvas | SUCCESS => {",
            "    \"ansible_facts\": { \"discovered_interpreter_python\": \"/usr/bin/python3\" },",
            "    \"changed\": false,",
            "    \"ping\": \"pong\"",
            "}"])


# ─── SLIDE 13: ANATOMY OF A PLAYBOOK ───────────────────────────────────────
s = content_slide("Anatomy of a Playbook — artists.yml",
                  "Source: docs.ansible.com/playbook_guide  |  /root/ansible/artists.yml on brush")

add_text(s, Inches(0.30), Inches(0.95), Inches(12.7), Inches(0.55),
         "A playbook is a YAML list of plays. Each play targets hosts and runs tasks. Each task calls a "
         "module — a small, idempotent unit of work that maps directly back to a Phase 1 command.",
         font=CALIBRI, size=12, color=TEXT)

code_block(s, Inches(0.30), Inches(1.65), Inches(7.50), Inches(5.40),
           "$ /root/ansible/artists.yml  (excerpt)",
           ["- name: Configure artists studio",
            "  hosts: canvas",
            "  become: true",
            "  vars:",
            "    artists: [vangogh, picasso, davinci,",
            "              monet, frida, warhol]",
            "",
            "  tasks:",
            "    - name: Create artists group",
            "      ansible.builtin.group:",
            "        name: artists",
            "        state: present",
            "",
            "    - name: Create artist users",
            "      ansible.builtin.user:",
            "        name: \"{{ item }}\"",
            "        group: artists",
            "        create_home: true",
            "      loop: \"{{ artists }}\""])

card(s, Inches(8.00), Inches(1.65), Inches(5.00), Inches(2.55),
     "Module → bash equivalent", accent=CYAN,
     items=["▸  group:  → groupadd",
            "▸  user:  → useradd -m -g …",
            "▸  file:  → mkdir + chown + chmod",
            "▸  copy:  → echo > … + chown + chmod",
            "▸  apt:    → apt install -y …",
            "▸  service: → systemctl reload …"], item_size=12)

card(s, Inches(8.00), Inches(4.30), Inches(5.00), Inches(2.75),
     "Why YAML over a bash script", accent=GREEN,
     items=["▸  Declarative — describe end state",
            "▸  Idempotent — only acts if needed",
            "▸  Reusable — swap inventory, same play",
            "▸  Auditable — diffs, dry-runs, recap"], item_size=12)


# ─── SLIDE 14: SANITY-CHECKING ─────────────────────────────────────────────
s = content_slide("Sanity-Checking Changes — --check, --diff, --step",
                  "Source: ansible-playbook(1)  |  docs.ansible.com/playbook_guide/playbooks_checkmode")

card(s, Inches(0.30), Inches(0.95), Inches(6.30), Inches(2.20),
     "--check  (plan mode)", accent=GREEN,
     body="Runs every task in 'would do' mode. Modules report what they would change without "
          "touching the host. Equivalent to terraform plan or apt --dry-run.",
     body_size=13)

card(s, Inches(6.70), Inches(0.95), Inches(6.30), Inches(2.20),
     "--diff  (show file changes)", accent=CYAN,
     body="Pairs with --check. For any file or template task, prints a unified diff of the before/after "
          "content — so you see the exact lines that would change.",
     body_size=13)

card(s, Inches(0.30), Inches(3.30), Inches(6.30), Inches(2.20),
     "--step  (pause between tasks)", accent=ORANGE,
     body="Prompts (y/n/c) before every task. Useful for slow, deliberate runs — and for narrating "
          "what is about to happen during a presentation.",
     body_size=13)

card(s, Inches(6.70), Inches(3.30), Inches(6.30), Inches(2.20),
     "--start-at-task / --tags / --limit", accent=RED,
     body="Cherry-pick: --start-at-task \"Name\" skips ahead, --tags users runs only tagged tasks, "
          "--limit canvas restricts to one host from the inventory.",
     body_size=13)

code_block(s, Inches(0.30), Inches(5.65), Inches(12.70), Inches(1.40),
           "$ three flavours of running the same playbook",
           ["ansible-playbook artists.yml --check --diff      # dry run, shows would-be diffs",
            "ansible-playbook artists.yml --step              # real, prompts before each task",
            "ansible-playbook artists.yml                      # real, full speed"])


# ─── SLIDE 15: IDEMPOTENCY ─────────────────────────────────────────────────
s = content_slide("Idempotency — the Property That Makes Ansible Safe to Re-Run",
                  "Source: docs.ansible.com  |  Wikipedia: idempotence (mathematics & systems)")

add_text(s, Inches(0.30), Inches(0.95), Inches(12.7), Inches(0.55),
         "Run the playbook a hundred times. The host ends up in the same state every single time. "
         "This is the property bash scripts don't have — and the reason ansible can be re-run safely.",
         font=CALIBRI, size=13, color=TEXT)

add_rect(s, Inches(0.30), Inches(1.65), Inches(6.30), Inches(3.40), CARD)
add_rect(s, Inches(0.30), Inches(1.65), Emu(54864), Inches(3.40), ORANGE)
add_text(s, Inches(0.45), Inches(1.73), Inches(6), Inches(0.40),
         "Bash on re-run", font=CALIBRI, size=14, bold=True, color=WHITE)
code_block(s, Inches(0.45), Inches(2.20), Inches(6.00), Inches(2.75),
           "# second run of the same script",
           [("$ groupadd painters", TEXT),
            ("groupadd: group 'painters' already exists", RED),
            ("$ useradd -m rembrandt", TEXT),
            ("useradd: user 'rembrandt' already exists", RED),
            ("$ echo $?", TEXT),
            ("9   # script aborts here unless guarded", RED)])

add_rect(s, Inches(6.70), Inches(1.65), Inches(6.30), Inches(3.40), CARD)
add_rect(s, Inches(6.70), Inches(1.65), Emu(54864), Inches(3.40), GREEN)
add_text(s, Inches(6.85), Inches(1.73), Inches(6), Inches(0.40),
         "Ansible on re-run", font=CALIBRI, size=14, bold=True, color=WHITE)
code_block(s, Inches(6.85), Inches(2.20), Inches(6.00), Inches(2.75),
           "$ ansible-playbook artists.yml  (second run)",
           [("TASK [Create artists group]", TEXT),
            ("ok: [canvas]", GREEN),
            ("TASK [Create artist users]", TEXT),
            ("ok: [canvas] => (item=vangogh)", GREEN),
            ("ok: [canvas] => (item=picasso)", GREEN),
            ("PLAY RECAP  canvas: ok=8  changed=0", GREEN)])

add_rect(s, Inches(0.30), Inches(5.25), Inches(12.70), Inches(1.80), CARD)
add_rect(s, Inches(0.30), Inches(5.25), Emu(54864), Inches(1.80), CYAN)
add_text(s, Inches(0.45), Inches(5.35), Inches(12.40), Inches(0.40),
         "Why this matters", font=CALIBRI, size=14, bold=True, color=WHITE)
add_multiline(s, Inches(0.45), Inches(5.75), Inches(12.40), Inches(1.20),
              ["▸  Run on schedule (cron, CI) without fear of duplicate side-effects.",
               "▸  Add a new server tomorrow — same playbook brings it to the same state without breaking the others.",
               "▸  changed: vs ok: in the recap is an audit trail. If something changed when it shouldn't have, it shows up there."],
              size=12, color=TEXT, line_spacing=1.2)


# ─── SLIDE 16: SERVICE INSTALL AT SCALE ────────────────────────────────────
s = content_slide("The Real Payoff — A Full Service Install in One Command",
                  "Source: /root/ansible/piwigo.yml on brush  |  piwigo.org/install")

add_text(s, Inches(0.30), Inches(0.95), Inches(12.7), Inches(0.55),
         "Users and permissions are warm-up. The point of the demo is what happens when this same "
         "approach is applied to a real service. piwigo.yml installs Apache + MariaDB + PHP + Piwigo + a vhost.",
         font=CALIBRI, size=12, color=TEXT)

card(s, Inches(0.30), Inches(1.65), Inches(6.30), Inches(4.20),
     "What this would be by hand (10 steps)", accent=ORANGE,
     items=["1.  apt install apache2 mariadb-server php …",
            "2.  CREATE DATABASE piwigo;",
            "3.  CREATE USER + GRANT ALL ON piwigo.*",
            "4.  curl piwigo.zip from upstream",
            "5.  unzip into /var/www/piwigo",
            "6.  chown -R www-data:www-data",
            "7.  write /etc/apache2/sites-available/piwigo.conf",
            "8.  a2enmod rewrite",
            "9.  a2ensite piwigo.conf",
            "10. systemctl reload apache2"], item_size=12)

card(s, Inches(6.70), Inches(1.65), Inches(6.30), Inches(4.20),
     "What it is in ansible", accent=GREEN,
     items=["▸  one ansible-playbook piwigo.yml",
            "▸  one YAML file describing the end state",
            "▸  reusable across N servers",
            "▸  --check shows the diff before any change",
            "▸  re-run is safe (idempotent)",
            "",
            "After it runs, browse to:",
            "    http://192.168.50.26/install.php",
            "and walk the install wizard — DB user 'piwigo',",
            "password 'BrushAndCanvas!2026'."], item_size=12)

code_block(s, Inches(0.30), Inches(6.05), Inches(12.70), Inches(1.00),
           "$ ten manual steps collapse to one command",
           ["ansible-playbook piwigo.yml"])


# ─── SLIDE 17: COMPARISON ──────────────────────────────────────────────────
s = content_slide("Side-by-Side — Manual vs Ansible",
                  "Italics = aspirational (the playbook does it internally; we don't run it by hand here)")

def trow(y, cells, fill, bold=False, font=CALIBRI, size=11, italic=False):
    cols = [(0.30, 4.30), (4.65, 4.65), (9.35, 3.65)]
    for (x, w), text in zip(cols, cells):
        add_rect(s, Inches(x), Inches(y), Inches(w), Inches(0.46), fill)
        tb = add_text(s, Inches(x + 0.10), Inches(y + 0.07), Inches(w - 0.20), Inches(0.36),
                      text, font=font, size=size, bold=bold, color=WHITE if bold else TEXT)
        if italic and tb.text_frame.paragraphs[0].runs:
            tb.text_frame.paragraphs[0].runs[0].font.italic = True

trow(0.95, ["Task", "Manual (bash)", "Ansible"], BAND, bold=True)

rows = [
    ("Create group",          "groupadd painters",                                     "group: name=artists"),
    ("Create 6 users",        "for u in …; do useradd -m -g painters \"$u\"; done",    "user: + loop:"),
    ("Directory + mode",      "mkdir; chown; chmod   (3 cmds)",                        "file: state=directory mode= …"),
    ("File + content + mode", "echo > …; chown; chmod   (3 cmds)",                     "copy: dest= content= mode= …"),
    ("Install 11 packages",   "apt install -y pkg1 pkg2 …",                            "apt: name=[…]"),
    ("Create DB + user",      "mysql <<SQL … SQL",                                      "mysql_db + mysql_user"),
    ("Reload service",        "systemctl reload apache2",                              "handler: service: state=reloaded"),
    ("Re-running",            "Error: group already exists",                           "ok  (idempotent)"),
    ("Across N servers",      "SSH into each one and repeat",                          "Add hostnames to inventory"),
]

y = 1.45
for i, (a, b, c) in enumerate(rows):
    fill = ALT_ROW if i % 2 == 0 else CARD
    italic = i in (4, 5, 6)
    trow(y, [a, b, c], fill, bold=False, font=CONSOLAS, size=11, italic=italic)
    y += 0.46

add_text(s, Inches(0.30), Inches(y + 0.10), Inches(12.70), Inches(0.35),
         "Italics = work the playbook does internally — not run by hand in this demo.",
         font=CONSOLAS, size=10, color=MUTED)


# ─── SLIDE 18: WHERE ANSIBLE FITS ──────────────────────────────────────────
s = content_slide("Where Ansible Fits — The Wider Infrastructure-as-Code Stack",
                  "Note: classmates are presenting Terraform separately. Ansible + Terraform are commonly paired.")

add_text(s, Inches(0.30), Inches(0.95), Inches(12.7), Inches(0.55),
         "Ansible is one tool in a layered stack. Each layer has a different job — provisioning, "
         "configuration, virtualisation, orchestration. They compose, not compete.",
         font=CALIBRI, size=12, color=TEXT)

card(s, Inches(0.30), Inches(1.70), Inches(6.30), Inches(2.30),
     "Provisioning  —  create the infrastructure", accent=CYAN,
     items=["▸  Terraform (HashiCorp HCL)  ← classmate's slot",
            "▸  Pulumi  (TypeScript / Python / Go)",
            "▸  CloudFormation, ARM, OpenTofu",
            "▸  Job: spin up VMs, networks, DNS, storage"], item_size=12)

card(s, Inches(6.70), Inches(1.70), Inches(6.30), Inches(2.30),
     "Configuration management  —  what we did today", accent=GREEN,
     items=["▸  Ansible       (this demo)",
            "▸  Puppet, Chef, Salt",
            "▸  cfengine (the original, 1993)",
            "▸  Job: bring an existing host to a desired state"], item_size=12)

card(s, Inches(0.30), Inches(4.10), Inches(6.30), Inches(2.30),
     "Virtualisation  —  what runs the VMs", accent=ORANGE,
     items=["▸  Proxmox VE  ← school stack",
            "▸  KVM / QEMU, ESXi, Hyper-V",
            "▸  Cloud equivalents: EC2, GCE, Azure VM",
            "▸  Job: provide the hosts the layers above target"], item_size=12)

card(s, Inches(6.70), Inches(4.10), Inches(6.30), Inches(2.30),
     "Orchestration  —  schedule containerised work", accent=RED,
     items=["▸  Kubernetes, Nomad",
            "▸  Docker Swarm, ECS",
            "▸  Job: place + scale containers across hosts",
            "▸  Adjacent to but distinct from the layers above"], item_size=12)

add_rect(s, Inches(0.30), Inches(6.55), Inches(12.70), Inches(0.55), CARD)
add_rect(s, Inches(0.30), Inches(6.55), Emu(54864), Inches(0.55), GREEN)
add_text(s, Inches(0.45), Inches(6.62), Inches(12.40), Inches(0.40),
         "Common pattern:  Terraform creates the VMs.  Ansible configures them.  Today we showed the second half.",
         font=CALIBRI, size=13, bold=True, color=WHITE)


# ─── SLIDE 19: CLOSING ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
slide_bg(s)
add_rect(s, 0, 0, Inches(0.50), prs.slide_height, GREEN)
add_text(s, Inches(0.70), Inches(1.00), Inches(12), Inches(2.40),
         "QUESTIONS?", font=IMPACT, size=110, bold=True, color=WHITE)
add_rect(s, Inches(0.70), Inches(3.30), Inches(9), Emu(54864), GREEN)
add_text(s, Inches(0.70), Inches(3.40), Inches(12), Inches(0.60),
         "USING LINUX — MANUALLY AND AUTOMATED",
         font=CONSOLAS, size=24, color=TEXT)
add_text(s, Inches(0.70), Inches(4.10), Inches(12), Inches(0.45),
         "Kernel  ●  Permissions  ●  FHS  ●  Inventory & playbooks  ●  Idempotency  ●  Service install at scale",
         font=CALIBRI, size=14, color=MUTED)
add_text(s, Inches(0.70), Inches(4.80), Inches(12), Inches(0.45),
         "kernel.org  ●  docs.ansible.com  ●  pathname.com/fhs  ●  piwigo.org",
         font=CONSOLAS, size=12, color=GREEN)
add_rect(s, 0, Inches(7.10), prs.slide_width, Inches(0.40), BAND)
add_text(s, Inches(0.30), Inches(7.13), Inches(12), Inches(0.28),
         "Career Practicum  ●  Finley Karras",
         font=CONSOLAS, size=9, color=MUTED)


prs.save(OUT)
print(f"wrote {OUT}")
print(f"slides: {len(prs.slides)}")
